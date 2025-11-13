"""
PRISM BRAIN AI V2 - Multi-Modal Research Analyzer
Complete implementation for standalone import
"""

import json
import re
import os
from datetime import datetime
from collections import defaultdict, Counter
from typing import Dict, List, Optional, Tuple
import hashlib
import requests

# Check for optional dependencies
try:
    import whisper
    WHISPER_AVAILABLE = True
    print("✓ Whisper available for audio processing")
except:
    WHISPER_AVAILABLE = False
    print("ℹ️  Whisper not installed (run: pip install openai-whisper)")

try:
    import PyPDF2
    PDF_AVAILABLE = True
    print("✓ PyPDF2 available for PDF processing")
except:
    PDF_AVAILABLE = False
    print("ℹ️  PyPDF2 not installed (run: pip install PyPDF2)")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
    print("✓ python-pptx available for PowerPoint processing")
except:
    PPTX_AVAILABLE = False
    print("ℹ️  python-pptx not installed (run: pip install python-pptx)")


class PRISMBrainV2:
    """
    PRISM Brain AI v2 - Multi-modal research synthesis

    Features:
    - FigJam: Full board analysis (sticky notes + arrows + diagrams)
    - Audio: Whisper transcription with tone analysis
    - Documents: PDF/PPT/DOCX analysis
    - Real-time synthesis and updates
    """

    def __init__(self, training_data=None, figjam_token=None):
        self.training_data = training_data
        self.figjam_token = figjam_token
        self.patterns = {}
        self.projects = {}

        # Initialize Whisper if available
        if WHISPER_AVAILABLE:
            try:
                print("🎙️  Loading Whisper model...")
                self.whisper_model = whisper.load_model("base")
                print("✓ Whisper model loaded")
            except Exception as e:
                print(f"⚠️  Whisper load error: {e}")
                self.whisper_model = None
        else:
            self.whisper_model = None

        # Learn from training data
        if training_data:
            self._initialize_from_training()

        print("\n🧠 PRISM Brain AI v2 initialized")
        print("   ✓ FigJam: URL input with full board analysis")
        print("   ✓ Audio: File upload with Whisper transcription")
        print("   ✓ Documents: PDF/PPT file upload")

    def _initialize_from_training(self):
        """Learn patterns from training data"""
        print("📚 Learning from training data...")
        all_notes = []
        for board in self.training_data[:100]:  # Sample for speed
            all_notes.extend(board['notes'])
        self.patterns['keywords'] = self._learn_keywords(all_notes)
        print(f"   ✓ Learned patterns from {len(all_notes):,} notes")

    def _learn_keywords(self, notes):
        """Learn keywords for content types"""
        keywords = defaultdict(set)
        for note in notes:
            content_type = note['true_type']
            words = note['content'].lower().split()
            keywords[content_type].update(words[:3])
        return {k: list(v)[:15] for k, v in keywords.items()}

    # =====================================================
    # PROJECT MANAGEMENT
    # =====================================================

    def create_project(self, project_name: str) -> str:
        """Create new PRISM project"""
        project_id = hashlib.md5(f"{project_name}{datetime.now()}".encode()).hexdigest()[:8]

        self.projects[project_id] = {
            'name': project_name,
            'created_at': datetime.now().isoformat(),
            'sources': [],
            'notes': [],
            'connections': [],
            'diagrams': [],
            'timeline': [],
            'contributors': {},
            'insights': {},
            'last_updated': datetime.now().isoformat()
        }

        print(f"✅ Created project: {project_name} (ID: {project_id})")
        return project_id

    # =====================================================
    # FIGJAM: URL UPLOAD WITH FULL ANALYSIS
    # =====================================================

    def ingest_figjam_url(self, project_id: str, figjam_url: str):
        """Ingest FigJam board from URL - analyzes everything"""
        print(f"\n📥 Ingesting FigJam board...")

        # Extract file key
        file_key = self._extract_figjam_key(figjam_url)
        if not file_key:
            return {'error': 'Invalid FigJam URL format'}

        if not self.figjam_token:
            return {'error': 'FigJam token not configured'}

        # Fetch board
        board_data = self._fetch_figjam_board(file_key)
        if not board_data:
            return {'error': 'Could not fetch board from API'}

        project = self.projects[project_id]
        board_name = board_data.get('name', 'Untitled Board')

        # Extract all elements
        sticky_notes = []
        arrows = []
        shapes = []

        def traverse(node, parent=None):
            node_type = node.get('type')

            if node_type == 'STICKY':
                sticky_notes.append({
                    'id': node.get('id'),
                    'content': node.get('characters', ''),
                    'color': self._map_color(node),
                    'author': node.get('lastModifier', {}).get('name', 'Unknown'),
                    'position': node.get('absoluteBoundingBox', {}),
                    'parent': parent
                })

            elif node_type == 'CONNECTOR':
                arrows.append({
                    'id': node.get('id'),
                    'from': node.get('connectorStart', {}).get('endpointNodeId'),
                    'to': node.get('connectorEnd', {}).get('endpointNodeId')
                })

            elif node_type in ['RECTANGLE', 'ELLIPSE', 'TEXT']:
                shapes.append({
                    'id': node.get('id'),
                    'type': node_type.lower(),
                    'content': node.get('characters', ''),
                    'position': node.get('absoluteBoundingBox', {})
                })

            if 'children' in node:
                for child in node['children']:
                    traverse(child, node.get('name') if node_type == 'FRAME' else parent)

        traverse(board_data.get('document', {}))

        print(f"   ✓ {len(sticky_notes)} sticky notes")
        print(f"   ✓ {len(arrows)} connections")
        print(f"   ✓ {len(shapes)} diagrams/shapes")

        # Analyze notes
        analyzed_notes = []
        for sticky in sticky_notes:
            analysis = self._analyze_content(sticky['content'], sticky['color'])

            note = {
                'id': sticky['id'],
                'source': 'figjam',
                'source_name': board_name,
                'content': sticky['content'],
                'color': sticky['color'],
                'predicted_type': analysis['predicted_type'],
                'confidence': analysis['confidence'],
                'contributor': sticky['author'],
                'created_at': datetime.now().isoformat(),
                'position': sticky['position'],
                'sentiment': self._detect_sentiment(sticky['content']),
                'priority': self._calc_priority(sticky['content'], analysis),
                'tags': self._extract_tags(sticky['content'])
            }
            analyzed_notes.append(note)

        # Store connections
        for arrow in arrows:
            project['connections'].append({
                'from_note': arrow['from'],
                'to_note': arrow['to'],
                'relationship': 'connects_to',
                'source': 'figjam'
            })

        project['diagrams'].extend(shapes)

        # Add to project
        project['sources'].append({
            'type': 'figjam',
            'name': board_name,
            'url': figjam_url,
            'added_at': datetime.now().isoformat(),
            'note_count': len(analyzed_notes),
            'connection_count': len(arrows),
            'diagram_count': len(shapes)
        })

        project['notes'].extend(analyzed_notes)
        project['last_updated'] = datetime.now().isoformat()

        self._update_timeline(project, analyzed_notes)
        self._update_contributors(project, analyzed_notes)

        print(f"   ✅ FigJam board ingested successfully")

        return {
            'success': True,
            'notes': len(analyzed_notes),
            'connections': len(arrows),
            'diagrams': len(shapes)
        }

    def _extract_figjam_key(self, url: str) -> Optional[str]:
        """Extract file key from FigJam URL"""
        match = re.search(r'/board/([a-zA-Z0-9_-]+)', url)
        if match:
            return match.group(1)
        match = re.search(r'/file/([a-zA-Z0-9_-]+)', url)
        if match:
            return match.group(1)
        return None

    def _fetch_figjam_board(self, file_key: str) -> Optional[Dict]:
        """Fetch from FigJam API"""
        url = f"https://api.figma.com/v1/files/{file_key}"
        headers = {"X-Figma-Token": self.figjam_token}

        try:
            response = requests.get(url, headers=headers, timeout=30)
            if response.status_code == 200:
                return response.json()
            else:
                print(f"   ❌ API error: {response.status_code}")
                return None
        except Exception as e:
            print(f"   ❌ Error: {e}")
            return None

    def _map_color(self, node):
        """Map RGB to color names"""
        fills = node.get('fills', [])
        if not fills or fills[0].get('type') != 'SOLID':
            return 'YELLOW'

        c = fills[0].get('color', {})
        r, g, b = c.get('r', 1), c.get('g', 1), c.get('b', 1)

        if r > 0.8 and g < 0.5 and b < 0.5:
            return 'RED'
        elif r > 0.8 and g > 0.5 and b < 0.3:
            return 'ORANGE'
        elif r > 0.8 and g > 0.8 and b < 0.5:
            return 'YELLOW'
        elif r < 0.5 and g > 0.7 and b < 0.5:
            return 'GREEN'
        elif r < 0.5 and g < 0.5 and b > 0.8:
            return 'BLUE'
        elif r > 0.5 and g < 0.5 and b > 0.7:
            return 'PURPLE'
        elif r > 0.8 and g < 0.5 and b > 0.6:
            return 'PINK'
        else:
            return 'GRAY'

    # =====================================================
    # AUDIO: FILE UPLOAD WITH WHISPER
    # =====================================================

    def ingest_audio_file(self, project_id: str, audio_path: str):
        """Ingest audio file with Whisper transcription"""
        print(f"\n🎙️  Processing audio file...")

        if not WHISPER_AVAILABLE or not self.whisper_model:
            return {'error': 'Whisper not available'}

        project = self.projects[project_id]
        file_name = os.path.basename(audio_path)

        print(f"   Transcribing: {file_name}")

        try:
            result = self.whisper_model.transcribe(
                audio_path,
                word_timestamps=True,
                verbose=False
            )
        except Exception as e:
            return {'error': f'Transcription failed: {e}'}

        segments = result['segments']
        print(f"   ✓ Transcribed {len(segments)} segments")

        # Extract insights
        analyzed_notes = []
        for i, seg in enumerate(segments):
            text = seg['text'].strip()
            if len(text) < 10:
                continue

            # Detect tone
            tone = self._analyze_tone(seg, text)

            # Extract key points
            points = self._extract_insights(text)

            for point in points:
                analysis = self._analyze_content(point, 'YELLOW')

                note = {
                    'id': f"audio_{file_name}_{i}_{len(analyzed_notes)}",
                    'source': 'audio',
                    'source_name': file_name,
                    'content': point,
                    'full_segment': text,
                    'predicted_type': analysis['predicted_type'],
                    'confidence': analysis['confidence'],
                    'contributor': 'Speaker',
                    'created_at': datetime.now().isoformat(),
                    'timestamp': f"{seg['start']:.1f}s",
                    'audio_tone': tone,
                    'sentiment': self._detect_sentiment(point),
                    'priority': self._calc_priority(point, analysis),
                    'tags': self._extract_tags(point)
                }
                analyzed_notes.append(note)

        project['sources'].append({
            'type': 'audio',
            'name': file_name,
            'added_at': datetime.now().isoformat(),
            'duration': f"{result.get('duration', 0):.1f}s",
            'note_count': len(analyzed_notes)
        })

        project['notes'].extend(analyzed_notes)
        project['last_updated'] = datetime.now().isoformat()

        self._update_timeline(project, analyzed_notes)
        self._update_contributors(project, analyzed_notes)

        print(f"   ✅ Extracted {len(analyzed_notes)} insights")

        return {
            'success': True,
            'notes': len(analyzed_notes),
            'duration': result.get('duration', 0)
        }

    def _analyze_tone(self, segment, text):
        """Detect speaker tone"""
        if '!' in text or text.isupper():
            return 'emphatic'
        elif '?' in text:
            return 'questioning'
        else:
            return 'neutral'

    def _extract_insights(self, text):
        """Extract key points from text"""
        points = []

        if any(w in text.lower() for w in ['problem', 'issue', 'difficult']):
            points.append(f"Pain point: {text}")
        elif '?' in text:
            points.append(f"Question: {text}")
        elif any(w in text.lower() for w in ['decided', 'agreed', 'will']):
            points.append(f"Decision: {text}")
        elif '"' in text:
            points.append(f"Quote: {text}")
        else:
            points.append(text)

        return points

    # =====================================================
    # DOCUMENTS: FILE UPLOAD
    # =====================================================

    def ingest_document_file(self, project_id: str, doc_path: str):
        """Ingest document file"""
        print(f"\n📄 Processing document...")

        file_name = os.path.basename(doc_path)
        ext = os.path.splitext(file_name)[1].lower()

        if ext == '.pdf':
            return self._ingest_pdf(project_id, doc_path)
        elif ext in ['.ppt', '.pptx']:
            return self._ingest_ppt(project_id, doc_path)
        else:
            return self._ingest_text(project_id, doc_path)

    def _ingest_pdf(self, project_id: str, pdf_path: str):
        """Extract from PDF"""
        if not PDF_AVAILABLE:
            return {'error': 'PyPDF2 not installed'}

        project = self.projects[project_id]
        file_name = os.path.basename(pdf_path)

        analyzed_notes = []

        try:
            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                print(f"   Processing {len(reader.pages)} pages")

                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    paras = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]

                    for para in paras:
                        analysis = self._analyze_content(para, 'YELLOW')

                        note = {
                            'id': f"pdf_{file_name}_p{page_num}_{len(analyzed_notes)}",
                            'source': 'pdf',
                            'source_name': file_name,
                            'content': para[:200],
                            'full_text': para,
                            'predicted_type': analysis['predicted_type'],
                            'confidence': analysis['confidence'],
                            'contributor': 'Author',
                            'created_at': datetime.now().isoformat(),
                            'page_number': page_num,
                            'sentiment': self._detect_sentiment(para),
                            'priority': self._calc_priority(para, analysis),
                            'tags': self._extract_tags(para)
                        }
                        analyzed_notes.append(note)
        except Exception as e:
            return {'error': f'PDF processing failed: {e}'}

        project['sources'].append({
            'type': 'pdf',
            'name': file_name,
            'added_at': datetime.now().isoformat(),
            'pages': len(reader.pages),
            'note_count': len(analyzed_notes)
        })

        project['notes'].extend(analyzed_notes)
        project['last_updated'] = datetime.now().isoformat()

        self._update_timeline(project, analyzed_notes)
        self._update_contributors(project, analyzed_notes)

        print(f"   ✅ Extracted {len(analyzed_notes)} insights")

        return {'success': True, 'notes': len(analyzed_notes)}

    def _ingest_ppt(self, project_id: str, ppt_path: str):
        """Extract from PowerPoint"""
        if not PPTX_AVAILABLE:
            return {'error': 'python-pptx not installed'}

        project = self.projects[project_id]
        file_name = os.path.basename(ppt_path)

        analyzed_notes = []

        try:
            prs = Presentation(ppt_path)
            print(f"   Processing {len(prs.slides)} slides")

            for slide_num, slide in enumerate(prs.slides, 1):
                text = ' '.join(shape.text for shape in slide.shapes if hasattr(shape, "text"))

                if len(text) > 20:
                    analysis = self._analyze_content(text, 'YELLOW')

                    note = {
                        'id': f"ppt_{file_name}_s{slide_num}",
                        'source': 'powerpoint',
                        'source_name': file_name,
                        'content': text[:200],
                        'full_text': text,
                        'predicted_type': analysis['predicted_type'],
                        'confidence': analysis['confidence'],
                        'contributor': 'Presenter',
                        'created_at': datetime.now().isoformat(),
                        'slide_number': slide_num,
                        'sentiment': self._detect_sentiment(text),
                        'priority': self._calc_priority(text, analysis),
                        'tags': self._extract_tags(text)
                    }
                    analyzed_notes.append(note)
        except Exception as e:
            return {'error': f'PPT processing failed: {e}'}

        project['sources'].append({
            'type': 'powerpoint',
            'name': file_name,
            'added_at': datetime.now().isoformat(),
            'slides': len(prs.slides),
            'note_count': len(analyzed_notes)
        })

        project['notes'].extend(analyzed_notes)
        self._update_timeline(project, analyzed_notes)

        print(f"   ✅ Extracted {len(analyzed_notes)} insights")

        return {'success': True, 'notes': len(analyzed_notes)}

    def _ingest_text(self, project_id: str, text_path: str):
        """Process text file"""
        project = self.projects[project_id]
        file_name = os.path.basename(text_path)

        with open(text_path, 'r') as f:
            text = f.read()

        analyzed_notes = []
        paras = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]

        for i, para in enumerate(paras):
            analysis = self._analyze_content(para, 'YELLOW')

            note = {
                'id': f"txt_{file_name}_{i}",
                'source': 'document',
                'source_name': file_name,
                'content': para[:200],
                'full_text': para,
                'predicted_type': analysis['predicted_type'],
                'confidence': analysis['confidence'],
                'contributor': 'Author',
                'created_at': datetime.now().isoformat(),
                'sentiment': self._detect_sentiment(para),
                'priority': self._calc_priority(para, analysis),
                'tags': self._extract_tags(para)
            }
            analyzed_notes.append(note)

        project['sources'].append({
            'type': 'document',
            'name': file_name,
            'added_at': datetime.now().isoformat(),
            'note_count': len(analyzed_notes)
        })

        project['notes'].extend(analyzed_notes)
        self._update_timeline(project, analyzed_notes)

        return {'success': True, 'notes': len(analyzed_notes)}

    # =====================================================
    # ANALYSIS FUNCTIONS
    # =====================================================

    def _analyze_content(self, content: str, color: str) -> Dict:
        """Analyze content type"""
        content_lower = content.lower()

        if '?' in content:
            return {'predicted_type': 'question', 'confidence': 0.8}
        elif '"' in content:
            return {'predicted_type': 'quote', 'confidence': 0.7}
        elif any(w in content_lower for w in ['problem', 'issue', 'error', 'broken']):
            return {'predicted_type': 'pain_point', 'confidence': 0.75}
        elif any(w in content_lower for w in ['love', 'great', 'awesome', 'excellent']):
            return {'predicted_type': 'positive', 'confidence': 0.7}
        elif any(w in content_lower for w in ['could', 'should', 'what if', 'idea']):
            return {'predicted_type': 'idea', 'confidence': 0.7}
        else:
            return {'predicted_type': 'neutral', 'confidence': 0.6}

    def _detect_sentiment(self, content: str) -> str:
        """Detect sentiment"""
        content_lower = content.lower()
        pos = sum(1 for w in ['love', 'great', 'good', 'excellent'] if w in content_lower)
        neg = sum(1 for w in ['hate', 'bad', 'terrible', 'broken'] if w in content_lower)
        return 'positive' if pos > neg else 'negative' if neg > pos else 'neutral'

    def _calc_priority(self, content: str, analysis: Dict) -> str:
        """Calculate priority"""
        if analysis['predicted_type'] == 'pain_point':
            return 'high'
        if any(w in content.lower() for w in ['critical', 'urgent', 'blocker']):
            return 'high'
        if analysis['predicted_type'] == 'neutral':
            return 'low'
        return 'medium'

    def _extract_tags(self, content: str) -> List[str]:
        """Extract tags"""
        tags = []
        content_lower = content.lower()

        tag_map = {
            'navigation': ['navigation', 'nav', 'menu'],
            'mobile': ['mobile', 'phone'],
            'performance': ['slow', 'fast', 'loading'],
            'accessibility': ['accessibility', 'a11y'],
            'search': ['search', 'find'],
            'error': ['error', 'bug', 'broken']
        }

        for tag, keywords in tag_map.items():
            if any(kw in content_lower for kw in keywords):
                tags.append(tag)

        return tags[:3]

    def _update_timeline(self, project: Dict, notes: List[Dict]):
        """Update timeline"""
        for note in notes:
            project['timeline'].append({
                'timestamp': note.get('created_at'),
                'contributor': note['contributor'],
                'content_preview': note['content'][:100],
                'note_id': note['id'],
                'source': note['source']
            })
        project['timeline'].sort(key=lambda x: x['timestamp'])

    def _update_contributors(self, project: Dict, notes: List[Dict]):
        """Update contributors"""
        for note in notes:
            contributor = note['contributor']
            if contributor not in project['contributors']:
                project['contributors'][contributor] = {
                    'total_contributions': 0,
                    'note_types': defaultdict(int)
                }

            project['contributors'][contributor]['total_contributions'] += 1
            project['contributors'][contributor]['note_types'][note['predicted_type']] += 1

    # =====================================================
    # SYNTHESIS
    # =====================================================

    def synthesize_project(self, project_id: str) -> Dict:
        """Generate project synthesis"""
        project = self.projects[project_id]
        notes = project['notes']

        by_type = defaultdict(list)
        by_priority = defaultdict(list)

        for note in notes:
            by_type[note['predicted_type']].append(note)
            by_priority[note['priority']].append(note)

        all_tags = []
        for note in notes:
            all_tags.extend(note.get('tags', []))

        tag_counts = Counter(all_tags)
        themes = [
            {'name': tag, 'frequency': count, 'percentage': (count/len(notes))*100}
            for tag, count in tag_counts.most_common(10)
        ]

        action_items = [
            {
                'content': n['content'], 
                'type': n['predicted_type'],
                'contributor': n['contributor'],
                'source': n['source_name']
            }
            for n in notes if n['priority'] == 'high'
        ][:20]

        sentiment_dist = Counter(n.get('sentiment', 'neutral') for n in notes)

        synthesis = {
            'project_name': project['name'],
            'last_updated': project['last_updated'],
            'total_notes': len(notes),
            'total_sources': len(project['sources']),
            'contributors': len(project['contributors']),
            'by_type': dict(by_type),
            'by_priority': dict(by_priority),
            'by_contributor': project['contributors'],
            'timeline': project['timeline'],
            'themes': themes,
            'action_items': action_items,
            'stats': {
                'sentiment_distribution': dict(sentiment_dist),
                'avg_confidence': sum(n.get('confidence', 0) for n in notes) / len(notes) if notes else 0
            }
        }

        project['insights'] = synthesis
        return synthesis

    def refresh_project(self, project_id: str):
        """Refresh analysis"""
        return self.synthesize_project(project_id)

print("✅ PRISM Brain AI v2 module loaded successfully!")
